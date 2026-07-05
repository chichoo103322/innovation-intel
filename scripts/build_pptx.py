#!/usr/bin/env python3
"""Build AI Agent presentation PPTX with Swiss design principles."""
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Design tokens ──────────────────────────────────────────────
INK = RGBColor(0x0D, 0x0D, 0x0D)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x00, 0x2F, 0xA7)  # IKB Blue
ACCENT_LIGHT = RGBColor(0xE8, 0xED, 0xF5)
GRAY = RGBColor(0x99, 0x99, 0x99)
HAIRLINE = RGBColor(0xDD, 0xDD, 0xDD)
INK_TINT = RGBColor(0x55, 0x55, 0x55)

W = Inches(13.333)  # 16:9
H = Inches(7.5)
FONT = "Inter"
FONT_ZH = "Microsoft YaHei UI"

prs = Presentation()
prs.slide_width = W
prs.slide_height = H


# ── Helpers ─────────────────────────────────────────────────────
def blank_slide():
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def add_bg(slide, color=PAPER):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill=None, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(0.5)
    return shape


def add_textbox(slide, left, top, width, height, text="", font_size=Pt(14),
                color=INK, bold=False, font_name=FONT, alignment=PP_ALIGN.LEFT,
                line_spacing=1.3):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.line_spacing = line_spacing
    return tf


def set_font(run, size=Pt(14), color=INK, bold=False, name=FONT):
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = name


def add_rich_textbox(slide, left, top, width, height):
    """Return text_frame for manual paragraph building."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def add_para(tf, text, size=Pt(14), color=INK, bold=False, name=FONT,
             align=PP_ALIGN.LEFT, spacing=1.3, space_after=Pt(4), first=False):
    if first:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = name
    p.alignment = align
    p.line_spacing = spacing
    p.space_after = space_after
    return p


def add_kicker(tf, text, first=True):
    return add_para(tf, text, size=Pt(10), color=ACCENT, bold=True,
                    name=FONT, spacing=1.2, first=first)


def add_hero_title(tf, text):
    return add_para(tf, text, size=Pt(44), color=INK, bold=False,
                    name=FONT, spacing=1.08, space_after=Pt(24))


def add_h2(tf, text):
    return add_para(tf, text, size=Pt(28), color=INK, bold=False,
                    name=FONT, spacing=1.15, space_after=Pt(16))


def add_body(tf, text, color=INK_TINT, size=Pt(13)):
    return add_para(tf, text, size=size, color=color, bold=False,
                    name=FONT, spacing=1.6, space_after=Pt(6))


def add_accent_block(slide, left, top, width, height, text, font_size=Pt(11)):
    shape = add_rect(slide, left, top, width, height, fill=ACCENT)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = PAPER
    p.font.bold = True
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER
    return tf


def add_card(slide, left, top, width, height, accent_top=True):
    """Add a card with optional blue top border."""
    card = add_rect(slide, left, top, width, height, fill=None, border=HAIRLINE)
    if accent_top:
        # accent line on top
        add_rect(slide, left, top, width, Pt(3), fill=ACCENT)
    return card


def bullet_list(tf, items, size=Pt(12), color=INK_TINT):
    for item in items:
        p = tf.add_paragraph()
        p.text = f"  {item}"
        p.font.size = size
        p.font.color.rgb = color
        p.font.name = FONT
        p.line_spacing = 1.5
        p.space_after = Pt(2)


MARGIN = Inches(1.2)
CONTENT_W = W - MARGIN * 2


# ═══════════════════════════════════════════════════════════════════
# S01 · COVER
# ═══════════════════════════════════════════════════════════════════
s = blank_slide()
add_bg(s, PAPER)
# Decorative accent square top-right
add_rect(s, W - Inches(2.5), Inches(1.2), Inches(1.4), Inches(1.4), fill=ACCENT)
# Kicker
tf = add_rich_textbox(s, MARGIN, Inches(1.8), Inches(8), Inches(0.4))
add_kicker(tf, "INNOVATION INTELLIGENCE · WORKFLOW AUTOMATION")
# Title
tf = add_rich_textbox(s, MARGIN, Inches(2.3), Inches(9), Inches(2.2))
add_hero_title(tf, "AI Agent")
add_para(tf, "如何搭建创新情报工作流", size=Pt(44), color=INK, bold=False,
         name=FONT_ZH, spacing=1.08)
# Subtitle
tf = add_rich_textbox(s, MARGIN, Inches(4.6), Inches(7), Inches(1))
add_body(tf, "一个人，三天，从零到全自动日报/周报/月报系统。", size=Pt(16))
add_body(tf, "不是魔法，是 Claude Code + Skill 生态的工程实践。", size=Pt(16))
# Meta
tf = add_rich_textbox(s, MARGIN, Inches(6.0), Inches(5), Inches(0.4))
add_para(tf, "2026.07 · 常州创新情报项目实战复盘", size=Pt(10),
         color=GRAY, name=FONT, spacing=1.2)

# ═══════════════════════════════════════════════════════════════════
# S02 · PROBLEM
# ═══════════════════════════════════════════════════════════════════
s = blank_slide()
add_bg(s, PAPER)
tf = add_rich_textbox(s, MARGIN, Inches(1.2), Inches(5), Inches(0.3))
add_kicker(tf, "THE PROBLEM")
tf = add_rich_textbox(s, MARGIN, Inches(1.6), Inches(5.5), Inches(1.2))
add_h2(tf, "每天搜情报，耗时 3 小时，还不全。")

tf = add_rich_textbox(s, MARGIN, Inches(3.0), Inches(5.5), Inches(0.8))
add_body(tf, "科技创新决策需要实时情报支撑。人工搜索覆盖不全、格式不统一、分发不及时。能不能让 AI 全自动搞定？")

# Stats row
for i, (num, label) in enumerate([("3h", "人工耗时/天"), ("4", "搜索维度"), ("0", "自动化程度")]):
    x = MARGIN + Inches(i * 2.2)
    tf = add_rich_textbox(s, x, Inches(4.0), Inches(2), Inches(1))
    add_para(tf, num, size=Pt(42), color=INK, bold=False, name=FONT, spacing=1.0)
    add_para(tf, label, size=Pt(10), color=GRAY, name=FONT, spacing=1.2)

# Goal card
card = add_rect(s, Inches(7.0), Inches(1.6), Inches(5.2), Inches(4.2), fill=ACCENT)
tf = add_rich_textbox(s, Inches(7.5), Inches(2.2), Inches(4.2), Inches(1.5))
add_para(tf, "目标：输入 API Key，输出桌面 PDF 报告，中间全自动。",
         size=Pt(18), color=PAPER, name=FONT_ZH, spacing=1.4)
# Tag pills
tags = ["联网搜索", "AI 分析", "HTML→PDF", "自动去重", "桌面分发", "定时调度"]
for i, tag in enumerate(tags):
    col = i % 3
    row = i // 3
    x = Inches(7.5) + Inches(col * 1.65)
    y = Inches(3.9) + Inches(row * 0.45)
    add_accent_block(s, x, y, Inches(1.4), Inches(0.32), tag, font_size=Pt(9))

# ═══════════════════════════════════════════════════════════════════
# S03 · THREE AGENTS
# ═══════════════════════════════════════════════════════════════════
s = blank_slide()
add_bg(s, INK)

agents = [
    ("Claude Code", "Anthropic · 2024", "终端里的 AI 编程伙伴。直接读写文件、执行命令、管理 Git。\n本项目的主力开发工具。", "本项目核心引擎", True),
    ("Cursor", "Anysphere · 2023", "AI-first 代码编辑器。逐文件精细编辑，\nTab 补全体验极好。适合前端和 UI 微调。", "UI 微调补充", False),
    ("GitHub Copilot", "Microsoft + OpenAI", "IDE 内嵌 AI 补全。日常编码加速。\nCoding Agent 模式可处理跨文件重构。", "日常编码加速", False),
]

tf = add_rich_textbox(s, MARGIN, Inches(0.9), Inches(8), Inches(0.8))
add_kicker(tf, "AI AGENT 三剑客", first=True)
add_kicker(tf, "", first=False)
add_para(tf, "三个 Agent，三种姿势", size=Pt(28), color=PAPER, name=FONT_ZH, spacing=1.15)

card_w = Inches(3.5)
card_h = Inches(4.6)
gap = Inches(0.4)
start_x = MARGIN

for i, (name, company, desc, tag, accent) in enumerate(agents):
    x = start_x + i * (card_w + gap)
    y = Inches(2.0)
    # Card border
    add_rect(s, x, y, card_w, card_h, fill=None,
             border=RGBColor(0x44, 0x44, 0x44) if not accent else ACCENT)
    # Icon square
    icon_fill = ACCENT if accent else RGBColor(0x33, 0x33, 0x33)
    add_rect(s, x + Inches(0.3), y + Inches(0.3), Inches(0.55), Inches(0.55), fill=icon_fill)
    # Company meta
    tf = add_rich_textbox(s, x + Inches(0.3), y + Inches(1.05), Inches(2.9), Inches(0.3))
    add_para(tf, company, size=Pt(8), color=GRAY, name=FONT, spacing=1.1)
    # Name
    tf = add_rich_textbox(s, x + Inches(0.3), y + Inches(1.35), Inches(2.9), Inches(0.5))
    add_para(tf, name, size=Pt(22), color=PAPER, name=FONT, spacing=1.15)
    # Desc
    tf = add_rich_textbox(s, x + Inches(0.3), y + Inches(2.0), Inches(2.9), Inches(1.4))
    add_body(tf, desc, color=RGBColor(0x99, 0x99, 0x99), size=Pt(11))
    # Tag
    tag_y = y + card_h - Inches(0.5)
    if accent:
        add_accent_block(s, x + Inches(0.3), tag_y, Inches(1.6), Inches(0.28), tag, Pt(8))
    else:
        tf = add_rich_textbox(s, x + Inches(0.3), tag_y, Inches(1.6), Inches(0.28))
        add_para(tf, tag, size=Pt(8), color=GRAY, name=FONT, spacing=1.1)

# ═══════════════════════════════════════════════════════════════════
# S04 · AGENT vs WEB CHAT
# ═══════════════════════════════════════════════════════════════════
s = blank_slide()
add_bg(s, INK)

tf = add_rich_textbox(s, MARGIN, Inches(0.7), Inches(8), Inches(0.7))
add_kicker(tf, "CORE DIFFERENCE")
add_h2(tf, "Agent ≠ 网页版 AI Chat")

tf = add_rich_textbox(s, MARGIN, Inches(1.7), Inches(8), Inches(0.4))
add_body(tf, "都是跟 AI 对话，但权限和能力天差地别", color=GRAY, size=Pt(13))

# Left: Web Chat
lx = MARGIN
ly = Inches(2.3)
lw = Inches(5.3)
lh = Inches(4.4)
add_rect(s, lx, ly, lw, lh, fill=None, border=RGBColor(0x33, 0x33, 0x33))
# Header
tf = add_rich_textbox(s, lx + Inches(0.3), ly + Inches(0.25), Inches(4.7), Inches(0.5))
add_para(tf, "网页版 AI Chat", size=Pt(18), color=GRAY, name=FONT_ZH, spacing=1.2)
add_para(tf, "chat.openai.com / guiaichat / Poe ...", size=Pt(9), color=RGBColor(0x66, 0x66, 0x66))
# Items
web_items = [
    "只能打字聊天 —— 你问，它答，纯文本",
    "复制粘贴工作流 —— 代码要手动拷到编辑器",
    "零系统权限 —— 不能读文件、不能写文件、不能执行命令",
    "无上下文记忆 —— 换个对话就失忆，无法持久工作",
    "无生态扩展 —— 不能装 Skill/MCP/插件",
]
tf = add_rich_textbox(s, lx + Inches(0.3), ly + Inches(1.1), Inches(4.7), Inches(2.4))
for item in web_items:
    add_para(tf, f"✕  {item}", size=Pt(12), color=RGBColor(0x88, 0x88, 0x88), spacing=1.6, space_after=Pt(4))
# Summary
add_rect(s, lx + Inches(0.3), ly + Inches(3.6), Inches(4.7), Inches(0.55), fill=RGBColor(0x1A, 0x1A, 0x1A))
tf = add_rich_textbox(s, lx + Inches(0.5), ly + Inches(3.7), Inches(4.3), Inches(0.4))
add_para(tf, "一句话：它只能告诉你答案，不能帮你做事", size=Pt(11), color=GRAY)

# Right: Agent
rx = Inches(7.0)
ry = Inches(2.3)
rw = Inches(5.3)
rh = Inches(4.4)
add_rect(s, rx, ry, rw, rh, fill=None, border=ACCENT)
tf = add_rich_textbox(s, rx + Inches(0.3), ry + Inches(0.25), Inches(4.7), Inches(0.5))
add_para(tf, "AI Agent（Claude Code）", size=Pt(18), color=PAPER, name=FONT_ZH, spacing=1.2)
add_para(tf, "终端里的 AI 工程师，有手有脚", size=Pt(9), color=RGBColor(0x66, 0x66, 0x99))
agent_items = [
    "文件系统读写 —— 直接创建/修改/删除项目文件",
    "Shell 命令执行 —— pip install、git commit、运行脚本",
    "联网搜索 + API 调用 —— 实时信息采集、第三方集成",
    "持久化项目记忆 —— CLAUDE.md/Memory，越用越懂你",
    "Skill + MCP 生态 —— 安装专业插件，无限扩展能力",
]
tf = add_rich_textbox(s, rx + Inches(0.3), ry + Inches(1.1), Inches(4.7), Inches(2.4))
for item in agent_items:
    add_para(tf, f"✓  {item}", size=Pt(12), color=RGBColor(0xBB, 0xBB, 0xBB), spacing=1.6, space_after=Pt(4))
# Summary
add_rect(s, rx + Inches(0.3), ry + Inches(3.6), Inches(4.7), Inches(0.55), fill=RGBColor(0x00, 0x1A, 0x55))
tf = add_rich_textbox(s, rx + Inches(0.5), ry + Inches(3.7), Inches(4.3), Inches(0.4))
add_para(tf, "一句话：它直接帮你把事做了，不用你动手", size=Pt(11), color=PAPER)

# ═══════════════════════════════════════════════════════════════════
# S05 · WHAT PERMISSIONS ENABLE
# ═══════════════════════════════════════════════════════════════════
s = blank_slide()
add_bg(s, PAPER)
tf = add_rich_textbox(s, MARGIN, Inches(0.8), Inches(8), Inches(0.7))
add_kicker(tf, "PERMISSION × CAPABILITY")
add_h2(tf, "有了这些权限，Agent 能做什么？")
tf = add_rich_textbox(s, MARGIN, Inches(1.7), Inches(8), Inches(0.4))
add_body(tf, "本项目实战。每一条都是 Claude Code 实际执行的。", size=Pt(13))

cards = [
    ("文件读写", "一次性创建 run_daily.py、generate_html_pdf.py 等 8 个核心脚本，自动写入几百行代码。不用你新建文件、不用你复制粘贴。"),
    ("Shell 执行", "pip3 install 装依赖，git init && git push 上传仓库，bash setup_scheduler.sh 部署定时任务。一条指令全自动。"),
    ("联网+API", "调用 DeepSeek API 联网搜索最新政策资讯，实时采集分析。搜索文档、查 GitHub、fetch 网页内容作为参考。"),
    ("Skill 扩展", "npx skills add 安装专业技能：anysearch 联网搜索、frontend-design 设计页面、guizang-ppt-skill 做演示文稿。装完即用。"),
]
card_w2 = Inches(2.65)
card_h2 = Inches(3.8)
gap2 = Inches(0.35)
for i, (title, desc) in enumerate(cards):
    x = MARGIN + i * (card_w2 + gap2)
    y = Inches(2.3)
    add_card(s, x, y, card_w2, card_h2)
    # Accent square icon
    add_rect(s, x + Inches(0.2), y + Inches(0.2), Inches(0.4), Inches(0.4), fill=ACCENT)
    tf = add_rich_textbox(s, x + Inches(0.2), y + Inches(0.75), Inches(2.25), Inches(0.3))
    add_para(tf, title, size=Pt(11), color=ACCENT, bold=True, name=FONT_ZH, spacing=1.1)
    tf = add_rich_textbox(s, x + Inches(0.2), y + Inches(1.15), Inches(2.25), Inches(2.4))
    add_body(tf, desc, size=Pt(11))

# ═══════════════════════════════════════════════════════════════════
# S06 · 7-STEP WORKFLOW TIMELINE
# ═══════════════════════════════════════════════════════════════════
s = blank_slide()
add_bg(s, PAPER)
tf = add_rich_textbox(s, MARGIN, Inches(0.8), Inches(8), Inches(0.7))
add_kicker(tf, "7-STEP WORKFLOW")
add_h2(tf, "一条命令，七个步骤，全自动流水线")

steps = ["环境配置", "API 接入", "Prompt 工程", "内容生成", "HTML→PDF", "桌面分发", "定时调度"]
desc_steps = ["Python + Chrome\n+ API Key", "DeepSeek SDK\n联网搜索", "系统提示词\n+ 维度设计", "JSON 输出\n+ 去重逻辑", "Chrome\n无头渲染", "自动存到\n桌面文件夹", "launchd/cron\n自动运行"]

# Horizontal line
add_rect(s, MARGIN, Inches(4.0), Inches(10.8), Pt(1), fill=ACCENT)

step_w = Inches(1.35)
gap_s = Inches(0.15)
for i, (step, desc) in enumerate(zip(steps, desc_steps)):
    x = MARGIN + i * (step_w + gap_s)
    # Dot
    dot_y = Inches(3.82)
    dot_fill = ACCENT if i == 6 else PAPER
    dot_border = ACCENT
    dot = add_rect(s, x + step_w/2 - Inches(0.08), dot_y, Inches(0.16), Inches(0.16),
                   fill=dot_fill, border=dot_border)
    # Label above
    tf = add_rich_textbox(s, x, Inches(2.2), step_w, Inches(1.0))
    add_para(tf, step, size=Pt(12), color=INK, bold=True, name=FONT_ZH, spacing=1.2, align=PP_ALIGN.CENTER)
    # Desc below
    tf = add_rich_textbox(s, x, Inches(4.3), step_w, Inches(1.2))
    add_para(tf, desc, size=Pt(10), color=GRAY, name=FONT_ZH, spacing=1.3, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# S07 · THREE LAYERS ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════
s = blank_slide()
add_bg(s, PAPER)
tf = add_rich_textbox(s, MARGIN, Inches(0.8), Inches(8), Inches(0.7))
add_kicker(tf, "ARCHITECTURE")
add_h2(tf, "三层架构：配置 → 引擎 → 输出")

layers = [
    ("01", "配置层 — config/settings.yaml",
     "API Key、模型选择、期号起始日期、桌面输出路径、调度频率。一个文件控制全局。", False),
    ("02", "引擎层 — scripts/ 核心脚本",
     "generate_html_pdf.py（所有报告生成+HTML→PDF）、dedup.py（去重）、distribute.py（分发）、setup_scheduler.sh（定时任务）", True),
    ("03", "输出层 — daily/ weekly/ monthly/",
     "日报（每工作日）、周报（每周五）、月报（每月末）→ HTML → PDF → 桌面「创新情报」文件夹", False),
]
ly = Inches(2.2)
for i, (num, title, desc, accent) in enumerate(layers):
    y = ly + i * Inches(1.5)
    fill = ACCENT if accent else None
    border = ACCENT if accent else HAIRLINE
    text_color = PAPER if accent else INK
    desc_color = RGBColor(0xCC, 0xCC, 0xDD) if accent else INK_TINT
    add_rect(s, MARGIN, y, CONTENT_W, Inches(1.2), fill=fill, border=border)
    # Number
    tf = add_rich_textbox(s, MARGIN + Inches(0.4), y + Inches(0.15), Inches(0.8), Inches(0.6))
    add_para(tf, num, size=Pt(32), color=ACCENT if not accent else PAPER, name=FONT, spacing=1.0)
    # Title
    tf = add_rich_textbox(s, MARGIN + Inches(1.3), y + Inches(0.2), Inches(9), Inches(0.4))
    add_para(tf, title, size=Pt(16), color=text_color, bold=True, name=FONT_ZH, spacing=1.2)
    # Desc
    tf = add_rich_textbox(s, MARGIN + Inches(1.3), y + Inches(0.6), Inches(9), Inches(0.5))
    add_para(tf, desc, size=Pt(12), color=desc_color, name=FONT_ZH, spacing=1.4)

# ═══════════════════════════════════════════════════════════════════
# S08 · API + PROMPT DETAIL
# ═══════════════════════════════════════════════════════════════════
s = blank_slide()
add_bg(s, PAPER)
tf = add_rich_textbox(s, MARGIN, Inches(0.8), Inches(8), Inches(0.7))
add_kicker(tf, "KEY STEP 1 & 2")
add_h2(tf, "API 接入 + Prompt 设计")

# Left card: API
add_card(s, MARGIN, Inches(2.2), Inches(5.3), Inches(4.2))
tf = add_rich_textbox(s, MARGIN + Inches(0.3), Inches(2.4), Inches(4.7), Inches(0.3))
add_para(tf, "接入 DeepSeek API", size=Pt(16), color=INK, bold=True, name=FONT_ZH, spacing=1.2)
tf = add_rich_textbox(s, MARGIN + Inches(0.3), Inches(2.8), Inches(4.7), Inches(0.5))
add_body(tf, "使用 OpenAI 兼容 SDK，只需改 base_url。\nenable_web_search: true 开启联网搜索。", size=Pt(12))
# Code block
code_bg = add_rect(s, MARGIN + Inches(0.3), Inches(3.5), Inches(4.7), Inches(2.5),
                   fill=RGBColor(0xF5, 0xF5, 0xF5))
tf = add_rich_textbox(s, MARGIN + Inches(0.5), Inches(3.6), Inches(4.3), Inches(2.3))
code_lines = [
    "client = OpenAI(",
    '  api_key="sk-...",',
    '  base_url="https://api.deepseek.com"',
    ")",
    "response = client.chat.completions.create(",
    '  model="deepseek-chat",',
    "  messages=[...],",
    '  extra_body={"enable_web_search": True}',
    ")",
]
for line in code_lines:
    add_para(tf, line, size=Pt(10), color=INK_TINT, name="Courier New", spacing=1.5, space_after=Pt(0))

# Right card: Prompt
add_card(s, Inches(7.0), Inches(2.2), Inches(5.3), Inches(4.2))
tf = add_rich_textbox(s, Inches(7.3), Inches(2.4), Inches(4.7), Inches(0.3))
add_para(tf, "System Prompt 设计要点", size=Pt(16), color=INK, bold=True, name=FONT_ZH, spacing=1.2)

prompt_items = [
    ("1. 角色设定", "资深科技创新情报分析师，服务于常州市科技创新决策"),
    ("2. 信源优先级", "gov.cn 政务官方 > 权威平台 > 媒体智库，每板块 ≥1 条来自政务官方"),
    ("3. 输出格式约束", "强制 JSON 输出 + 严格字数限制 + URL 必须"),
    ("4. 质量自检", "在 Prompt 中嵌入筛选标准和自检清单，让模型自己把关"),
]
for i, (title, desc) in enumerate(prompt_items):
    y = Inches(2.85) + i * Inches(0.82)
    tf = add_rich_textbox(s, Inches(7.3), y, Inches(4.7), Inches(0.25))
    add_para(tf, title, size=Pt(11), color=ACCENT, bold=True, name=FONT_ZH, spacing=1.1)
    tf = add_rich_textbox(s, Inches(7.3), y + Inches(0.22), Inches(4.7), Inches(0.5))
    add_para(tf, desc, size=Pt(10), color=INK_TINT, name=FONT_ZH, spacing=1.3)

# ═══════════════════════════════════════════════════════════════════
# S09 · BEFORE vs AFTER
# ═══════════════════════════════════════════════════════════════════
s = blank_slide()
add_bg(s, INK)
tf = add_rich_textbox(s, MARGIN, Inches(0.7), Inches(8), Inches(0.7))
add_kicker(tf, "BEFORE / AFTER")
add_h2(tf, "传统开发 vs Agent 协作")

# Left: old
lx = MARGIN
ly = Inches(2.0)
lw = Inches(5.3)
lh = Inches(4.6)
add_rect(s, lx, ly, lw, lh, fill=None, border=RGBColor(0x33, 0x33, 0x33))
tf = add_rich_textbox(s, lx + Inches(0.3), ly + Inches(0.2), Inches(4.7), Inches(0.35))
add_para(tf, "传统方式", size=Pt(16), color=GRAY, name=FONT_ZH, spacing=1.2)
old_items = [
    "学习 Python + 各种库：2 周",
    "研究 DeepSeek API 文档：2 天",
    "设计 Prompt + 反复调试：1 周",
    "HTML→PDF 方案选型 + 踩坑：3 天",
    "部署 + 定时任务调试：2 天",
    "写部署文档：1 天",
]
tf = add_rich_textbox(s, lx + Inches(0.3), ly + Inches(0.7), Inches(4.7), Inches(2.6))
for item in old_items:
    add_para(tf, f"✕  {item}", size=Pt(13), color=RGBColor(0x88, 0x88, 0x88), spacing=1.8, space_after=Pt(2))
tf = add_rich_textbox(s, lx + Inches(0.3), ly + Inches(3.5), Inches(4.7), Inches(0.8))
add_para(tf, "≈ 2-3 个月", size=Pt(36), color=RGBColor(0x66, 0x66, 0x66), spacing=1.0)

# Right: agent
rx = Inches(7.0)
ry = Inches(2.0)
rw = Inches(5.3)
rh = Inches(4.6)
add_rect(s, rx, ry, rw, rh, fill=None, border=ACCENT)
tf = add_rich_textbox(s, rx + Inches(0.3), ry + Inches(0.2), Inches(4.7), Inches(0.35))
add_para(tf, "Claude Code 方式", size=Pt(16), color=PAPER, name=FONT_ZH, spacing=1.2)
cc_items = [
    "\"帮我写一个调用 DeepSeek 的脚本\"：10 分钟",
    "\"设计一个日报 Prompt，要求……\"：30 分钟",
    "\"把 HTML 转成 PDF，用 Chrome headless\"：5 分钟",
    "\"给日报/周报/月报统一用这个管道\"：20 分钟",
    "\"用 launchd 每 30 分钟跑一次\"：10 分钟",
    "\"写一份部署手册，给小白看的\"：5 分钟",
]
tf = add_rich_textbox(s, rx + Inches(0.3), ry + Inches(0.7), Inches(4.7), Inches(2.6))
for item in cc_items:
    add_para(tf, f"✓  {item}", size=Pt(13), color=RGBColor(0xBB, 0xBB, 0xBB), spacing=1.8, space_after=Pt(2))
tf = add_rich_textbox(s, rx + Inches(0.3), ry + Inches(3.5), Inches(4.7), Inches(0.8))
add_para(tf, "≈ 2-3 天", size=Pt(36), color=ACCENT, spacing=1.0)

# ═══════════════════════════════════════════════════════════════════
# S10 · SKILL ECOSYSTEM
# ═══════════════════════════════════════════════════════════════════
s = blank_slide()
add_bg(s, PAPER)
tf = add_rich_textbox(s, MARGIN, Inches(0.7), Inches(8), Inches(0.7))
add_kicker(tf, "KEY STEP 3")
add_h2(tf, "Skill 生态：装一个 Skill，多一种超能力")
tf = add_rich_textbox(s, MARGIN, Inches(1.5), Inches(9), Inches(0.4))
add_body(tf, "Claude Code 的 Skill 系统 = 给 AI 装专业插件。每个 Skill 是一套专业知识和工具链。", size=Pt(13))

skills = [
    ("01", "anysearch", "统一实时搜索。联网查资料、验证事实、提取网页全文。本项目用于辅助验证 AI 生成内容的准确性。"),
    ("02", "frontend-design", "前端界面设计。避免 AI 美学——每次输出独特、生产级质量的 UI。本项目用于 HTML 报告模板设计。"),
    ("03", "guizang-ppt-skill", "瑞士国际主义风网页 PPT。本项目用于生成这份演示文稿本身——你没看错，这份 PPT 就是用它的设计原则做的。"),
    ("04", "humanize-ppt", "PPT 叙事结构化引擎。AST（受众状态转移）理论，每页翻动都推进叙事。本项目用于优化演示逻辑。"),
]
for i, (num, name, desc) in enumerate(skills):
    x = MARGIN + i * Inches(2.8)
    y = Inches(2.2)
    add_card(s, x, y, Inches(2.55), Inches(3.5))
    tf = add_rich_textbox(s, x + Inches(0.2), y + Inches(0.15), Inches(0.5), Inches(0.5))
    add_para(tf, num, size=Pt(28), color=ACCENT, spacing=1.0)
    tf = add_rich_textbox(s, x + Inches(0.2), y + Inches(0.7), Inches(2.15), Inches(0.3))
    add_para(tf, name, size=Pt(14), color=INK, bold=True, name=FONT, spacing=1.1)
    tf = add_rich_textbox(s, x + Inches(0.2), y + Inches(1.15), Inches(2.15), Inches(2.1))
    add_body(tf, desc, size=Pt(10))

# Bottom insight
add_rect(s, MARGIN, Inches(5.95), CONTENT_W, Inches(0.6), fill=ACCENT_LIGHT)
tf = add_rich_textbox(s, MARGIN + Inches(0.3), Inches(6.0), CONTENT_W - Inches(0.6), Inches(0.5))
add_para(tf, "核心经验：不是「装越多越好」，而是「在合适的环节用合适的 Skill」。搜索用 anysearch，设计用 frontend-design，PPT 用 guizang——各司其职。",
         size=Pt(10), color=INK_TINT, spacing=1.4)

# ═══════════════════════════════════════════════════════════════════
# S11 · DAILY LOOP
# ═══════════════════════════════════════════════════════════════════
s = blank_slide()
add_bg(s, PAPER)
tf = add_rich_textbox(s, MARGIN, Inches(0.8), Inches(8), Inches(0.7))
add_kicker(tf, "AUTOMATION LOOP")
add_h2(tf, "日报自动生成闭环")
tf = add_rich_textbox(s, MARGIN, Inches(1.6), Inches(8), Inches(0.4))
add_body(tf, "python3 run_daily.py  一条命令，七个环节全自动", size=Pt(13))

loop_steps = [
    ("联网搜索", "DeepSeek API\n4维度×2-3条"),
    ("AI 分析", "JSON 结构化\n摘要+创新洞察"),
    ("去重", "标题+URL\n历史记录比对"),
    ("HTML 生成", "封面+4板块\n统一设计语言"),
    ("HTML→PDF", "Chrome headless\n--print-to-pdf"),
    ("桌面分发", "~/Desktop/创新情报/\n日报/周报/月报/"),
    ("定时触发", "launchd 每30分钟\n已生成则跳过"),
]
step_w3 = Inches(1.45)
gap3 = Inches(0.12)
# Connecting line
add_rect(s, MARGIN + Inches(0.6), Inches(4.0), Inches(10.4), Pt(1), fill=ACCENT_LIGHT)
for i, (step, desc) in enumerate(loop_steps):
    x = MARGIN + i * (step_w3 + gap3)
    # Circle
    add_rect(s, x + step_w3/2 - Inches(0.11), Inches(3.82), Inches(0.22), Inches(0.22),
             fill=ACCENT if i == 6 else PAPER, border=ACCENT)
    # Title
    tf = add_rich_textbox(s, x, Inches(2.2), step_w3, Inches(0.8))
    add_para(tf, step, size=Pt(12), color=INK, bold=True, name=FONT_ZH, spacing=1.2, align=PP_ALIGN.CENTER)
    # Desc
    tf = add_rich_textbox(s, x, Inches(4.3), step_w3, Inches(1.2))
    add_para(tf, desc, size=Pt(10), color=GRAY, name=FONT_ZH, spacing=1.3, align=PP_ALIGN.CENTER)

# Arrow between each
for i in range(6):
    x = MARGIN + (i+1) * (step_w3 + gap3) - gap3/2
    tf = add_rich_textbox(s, x - Inches(0.05), Inches(3.6), Inches(0.2), Inches(0.4))
    add_para(tf, "→", size=Pt(14), color=ACCENT, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# S12 · RESULTS KPI
# ═══════════════════════════════════════════════════════════════════
s = blank_slide()
add_bg(s, PAPER)
tf = add_rich_textbox(s, MARGIN, Inches(0.8), Inches(8), Inches(0.7))
add_kicker(tf, "RESULTS")
add_h2(tf, "项目成果速览")

kpis = [
    ("3", "报告类型", "日报·周报·月报", 18),
    ("8-12", "条/期", "4维度×2-3条", 24),
    ("3min", "全流程耗时", "搜索→PDF", 30),
    ("0", "人工干预", "全自动运行", 36),
    ("∞", "可扩展", "加维度=加Prompt", 42),
]
bar_w = Inches(1.6)
bar_gap = Inches(0.6)
for i, (num, label, sub, h_pct) in enumerate(kpis):
    x = MARGIN + Inches(1.5) + i * (bar_w + bar_gap)
    # Value
    tf = add_rich_textbox(s, x, Inches(1.8), bar_w, Inches(1.0))
    add_para(tf, num, size=Pt(36), color=ACCENT if i >= 3 else INK,
             name=FONT, spacing=1.0, align=PP_ALIGN.CENTER)
    # Bar
    bar_h = Inches(h_pct * 0.06)
    bar_y = Inches(6.2) - bar_h
    bar_fill = ACCENT if i >= 3 else ACCENT_LIGHT
    add_rect(s, x + Inches(0.3), bar_y, Inches(1.0), bar_h, fill=bar_fill)
    # Label
    tf = add_rich_textbox(s, x, Inches(6.3), bar_w, Inches(0.5))
    add_para(tf, label, size=Pt(12), color=INK, bold=True, name=FONT_ZH, spacing=1.2, align=PP_ALIGN.CENTER)
    tf = add_rich_textbox(s, x, Inches(6.65), bar_w, Inches(0.35))
    add_para(tf, sub, size=Pt(9), color=GRAY, name=FONT_ZH, spacing=1.1, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# S13 · KEY INSIGHTS
# ═══════════════════════════════════════════════════════════════════
s = blank_slide()
add_bg(s, INK)
tf = add_rich_textbox(s, MARGIN, Inches(0.9), Inches(8), Inches(0.7))
add_kicker(tf, "CORE INSIGHT")
add_h2(tf, "不是 AI 替你写代码，是你指挥 AI 写代码。")

insights = [
    ("01", "Prompt 是新的编程语言",
     "写好 System Prompt 比写代码更重要。约束力、格式要求、自检清单，都在 Prompt 里定义清楚。"),
    ("02", "Skill = 专业外包",
     "你不会设计？装 frontend-design。不会做 PPT？装 guizang-ppt-skill。把专业领域外包给专业 Skill。"),
    ("03", "迭代，不是一次到位",
     "第一版总是有问题的。和 AI 对话式迭代——「封面太AI了，换一种风格」——每次 10 秒，改到满意。"),
]
for i, (num, title, desc) in enumerate(insights):
    y = Inches(2.3) + i * Inches(1.5)
    # Number badge
    add_rect(s, MARGIN, y + Inches(0.1), Inches(0.42), Inches(0.32), fill=ACCENT)
    tf = add_rich_textbox(s, MARGIN + Inches(0.02), y + Inches(0.1), Inches(0.38), Inches(0.32))
    add_para(tf, num, size=Pt(9), color=PAPER, bold=True, align=PP_ALIGN.CENTER)
    # Title
    tf = add_rich_textbox(s, MARGIN + Inches(0.7), y, Inches(5), Inches(0.4))
    add_para(tf, title, size=Pt(18), color=PAPER, bold=True, name=FONT_ZH, spacing=1.2)
    # Desc
    tf = add_rich_textbox(s, MARGIN + Inches(0.7), y + Inches(0.45), Inches(8), Inches(0.6))
    add_body(tf, desc, color=RGBColor(0x99, 0x99, 0x99), size=Pt(13))

# ═══════════════════════════════════════════════════════════════════
# S14 · ROADMAP
# ═══════════════════════════════════════════════════════════════════
s = blank_slide()
add_bg(s, PAPER)
tf = add_rich_textbox(s, MARGIN, Inches(0.8), Inches(8), Inches(0.7))
add_kicker(tf, "NEXT STEPS")
add_h2(tf, "三阶段行动路线")

phases = [
    ("PHASE 1 · 已完成", "自动化情报管线",
     "日报/周报/月报全自动生成，HTML→PDF 统一管道，去重+桌面分发+定时调度全部就绪。"),
    ("PHASE 2 · 进行中", "信源扩展 + 专题深挖",
     "接入更多垂直信源（专利数据库、学术论文、投融资数据），增加专题深度分析模块。"),
    ("PHASE 3 · 规划中", "交互式情报仪表盘",
     "Web 端可视化面板，支持关键词订阅、趋势图表、竞合关系图谱、邮件/飞书推送。"),
    ("PHASE X · 探索", "AI Agent 多机协作",
     "多 Agent 分工：采集 Agent + 分析 Agent + 排版 Agent，消息队列串联，真正「AI 团队」。"),
]
for i, (phase, title, desc) in enumerate(phases):
    x = MARGIN + i * Inches(2.8)
    y = Inches(2.2)
    add_card(s, x, y, Inches(2.55), Inches(4.2))
    tf = add_rich_textbox(s, x + Inches(0.2), y + Inches(0.15), Inches(2.15), Inches(0.25))
    add_para(tf, phase, size=Pt(8), color=ACCENT, bold=True, name=FONT, spacing=1.1)
    tf = add_rich_textbox(s, x + Inches(0.2), y + Inches(0.5), Inches(2.15), Inches(0.8))
    add_para(tf, title, size=Pt(16), color=INK, bold=True, name=FONT_ZH, spacing=1.2)
    tf = add_rich_textbox(s, x + Inches(0.2), y + Inches(1.4), Inches(2.15), Inches(2.5))
    add_body(tf, desc, size=Pt(11))

# ═══════════════════════════════════════════════════════════════════
# S15 · CLOSING
# ═══════════════════════════════════════════════════════════════════
s = blank_slide()
add_bg(s, ACCENT)
# Decorative square
add_rect(s, W - Inches(2.2), Inches(1.2), Inches(1.2), Inches(1.2),
         fill=None, border=RGBColor(0x66, 0x88, 0xCC))
tf = add_rich_textbox(s, MARGIN, Inches(2.0), Inches(9), Inches(1.5))
add_para(tf, "一个人 + Claude Code", size=Pt(48), color=PAPER,
         name=FONT_ZH, spacing=1.1)
add_para(tf, "= 一支技术团队", size=Pt(48), color=PAPER,
         name=FONT_ZH, spacing=1.1)
tf = add_rich_textbox(s, MARGIN, Inches(4.0), Inches(7), Inches(0.8))
add_para(tf, "AI Agent 不会取代你的工作。", size=Pt(18), color=RGBColor(0x88, 0x99, 0xCC),
         name=FONT_ZH, spacing=1.5)
add_para(tf, "但它会取代「不会用 AI Agent 的人」的工作。", size=Pt(18), color=RGBColor(0x88, 0x99, 0xCC),
         name=FONT_ZH, spacing=1.5)
# CTA pills
ctas = ["开始用起来", "从一个小项目开始", "装一个 Skill 试试"]
for i, cta in enumerate(ctas):
    x = MARGIN + Inches(i * 2.2)
    add_rect(s, x, Inches(5.2), Inches(1.9), Inches(0.45),
             fill=None, border=RGBColor(0x44, 0x55, 0x88))
    tf = add_rich_textbox(s, x, Inches(5.25), Inches(1.9), Inches(0.4))
    add_para(tf, cta, size=Pt(12), color=PAPER, align=PP_ALIGN.CENTER, spacing=1.2)

# ═══════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════
output = Path("/Users/jzxzhou/Desktop/创新情报/AI_Agent_创新情报工作流.pptx")
prs.save(str(output))
print(f"Saved: {output}")
print(f"Slides: {len(prs.slides)}")
